# -*- coding: utf-8 -*-
"""
build_sample_deck.py — design-system.md 규격대로 5장짜리 덱을 만든다.

수강생이 실제로 배우는 것은 이 파일이다.
Claude 에게 덱을 시킬 때도 결국 이런 코드가 만들어진다. 여기서 볼 것은 세 가지다.

  1. 존 좌표를 상수로 박아 두고 모든 슬라이드가 같은 값을 쓴다  → ZONE
  2. 색·서체를 상수로 모아 둔다                                  → BRAND (brand.config.md 와 1:1)
  3. 폭을 눈대중하지 않고 계산해서 쓴다                           → kpi_widths()

실행:  python3 build_sample_deck.py
필요:  pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import copy, os

# ────────────────────────────────────────────────────────────
# brand.config.md 와 1:1로 대응한다. 브랜드를 바꾸면 여기만 고친다.
# ────────────────────────────────────────────────────────────
BRAND = dict(
    name            = "MyBrand",
    font            = "Pretendard",
    primary         = RGBColor(0x14, 0x56, 0xf0),
    primary_light   = RGBColor(0x3b, 0x82, 0xf6),
    primary_lighter = RGBColor(0x60, 0xa5, 0xfa),
    text_main       = RGBColor(0x22, 0x22, 0x22),
    text_sub        = RGBColor(0x45, 0x51, 0x5e),
    text_muted      = RGBColor(0x8e, 0x8e, 0x93),
    white           = RGBColor(0xff, 0xff, 0xff),
    dark_surface    = RGBColor(0x18, 0x1e, 0x25),
    border          = RGBColor(0xe5, 0xe7, 0xeb),
    divider         = RGBColor(0xf2, 0xf3, 0xf5),
)
LOGO = os.path.join(os.path.dirname(__file__), "..", "assets", "logo-placeholder.png")

# ── 존 좌표 (design-system.md §0). 이 값은 슬라이드마다 바뀌지 않는다 ──
ZONE = dict(
    margin      = 0.5,
    header_top  = 0.40,
    logo_top    = 0.44,
    title_top   = 1.00, title_h = 0.75,
    sub_top     = 1.75, sub_h   = 0.40,
    body_top    = 2.39, body_bottom = 6.85,   # 높이 4.46"
    footer_top  = 7.05,
)
CONTENT_W = 12.333          # 13.333 − 좌우 0.5 × 2
GUTTER    = 0.2

def kpi_widths(n):
    """카드 n장을 내용 폭에 정확히 채우는 카드 폭. 눈대중하지 않는다.
    4장 → 2.93",  3장 → 3.98"  (design-system.md §5 패턴 A와 같은 값)"""
    return round((CONTENT_W - GUTTER * (n - 1)) / n, 2)

# ────────────────────────────────────────────────────────────
# 도우미
# ────────────────────────────────────────────────────────────
def textbox(slide, x, y, w, h, text, size, bold=False, color=None,
            align=PP_ALIGN.LEFT, line=1.3, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run(); r.text = ln
        r.font.name = BRAND["font"]; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color or BRAND["text_main"]
    return tb

def card(slide, x, y, w, h, fill=None, radius=0.06, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill or BRAND["white"]
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def gradient_card(slide, x, y, w, h):
    """Hero Gradient — 135°, Primary → PrimaryLight → PrimaryLighter.
    덱 전체에서 3개까지만. 이 덱은 표지 1개만 쓴다."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.gradient()
    stops = sh.fill.gradient_stops
    stops[0].color.rgb = BRAND["primary"];         stops[0].position = 0.0
    stops[1].color.rgb = BRAND["primary_lighter"]; stops[1].position = 1.0
    sh.fill.gradient_angle = 135.0
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def frame(slide, chapter, page, total, source="", dark=False):
    """모든 슬라이드에 같은 좌표로 들어가는 틀. 여기가 이 시스템의 핵심이다."""
    muted = RGBColor(0x9a, 0xa3, 0xad) if dark else BRAND["text_muted"]
    textbox(slide, ZONE["margin"], ZONE["header_top"], 6, 0.3,
            chapter, 12, bold=True, color=muted)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(13.333 - ZONE["margin"] - 1.22),
                                 Inches(ZONE["logo_top"]), Inches(1.22), Inches(0.24))
    textbox(slide, ZONE["margin"], ZONE["footer_top"], 3, 0.25,
            f"{page} / {total}", 10, color=muted)
    if source:
        textbox(slide, 13.333 - ZONE["margin"] - 7, ZONE["footer_top"], 7, 0.25,
                source, 9, color=muted, align=PP_ALIGN.RIGHT)

def head(slide, title, sub, dark=False):
    """제목은 한 줄로 쓴다.
    ★ 실측 함정 — 제목 존은 1.00″–1.75″ 로 0.75″ 뿐이다.
      36pt × 행간 1.20 × 2줄 = 86.4pt = 1.20″ 라서 두 줄을 쓰면 부제 존을 덮는다.
      두 줄이 필요하면 제목을 줄이거나, 제목 존을 아래로 넓히고 부제를 함께 내려야 한다."""
    textbox(slide, ZONE["margin"], ZONE["title_top"], CONTENT_W, ZONE["title_h"],
            title, 36, bold=True,
            color=BRAND["white"] if dark else BRAND["text_main"], line=1.20)
    textbox(slide, ZONE["margin"], ZONE["sub_top"], CONTENT_W, ZONE["sub_h"],
            sub, 16, color=RGBColor(0xc9,0xd0,0xd8) if dark else BRAND["text_sub"],
            line=1.45)

def won(n):
    """금액은 통화 최소 단위 + 세 자리 콤마. 만/억/조로 줄이지 않는다."""
    return f"{n:,}원"

# ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 5
CHAPTER = "SLIDE DESIGN SYSTEM"

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

# ── 1. 표지 — Hero Gradient 1개 (덱 예산 3개 중 1개 사용) ──
s = prs.slides.add_slide(BLANK); bg(s, BRAND["white"])
frame(s, CHAPTER, 1, TOTAL)
head(s, "좌표를 고정하면 덱이 흔들리지 않습니다",
     "이 다섯 장은 전부 같은 존 좌표를 씁니다. 바뀌는 것은 본문뿐입니다")
gradient_card(s, ZONE["margin"], ZONE["body_top"], CONTENT_W, 3.2)
textbox(s, ZONE["margin"] + 0.5, ZONE["body_top"] + 0.7, 8, 1.4,
        "4.46″", 60, bold=True, color=BRAND["white"], line=1.10)
textbox(s, ZONE["margin"] + 0.5, ZONE["body_top"] + 2.1, 10, 0.5,
        "모든 슬라이드가 공유하는 본문 상자 높이 (2.39″ – 6.85″)",
        12, color=RGBColor(0xdc,0xe8,0xff))
textbox(s, ZONE["margin"], ZONE["body_top"] + 3.4, CONTENT_W, 0.9,
        "이 덱은 design-system.md 규격을 그대로 따른 예시입니다.\n"
        "로고 자리에는 교체용 자리표시 이미지가 들어가 있습니다.",
        13, color=BRAND["text_sub"], line=1.5)

# ── 2. 패턴 A — KPI 4장 + 가로 막대 차트 ──
s = prs.slides.add_slide(BLANK); bg(s, BRAND["white"])
frame(s, CHAPTER, 2, TOTAL, "가상 데이터 · 서식 예시용")
head(s, "패턴 A — KPI 띠와 차트", "카드 폭은 눈대중이 아니라 계산해서 넣습니다")
w4 = kpi_widths(4)                                   # 2.93"
kpis = [("12", "본문 그리드 열 수"), ("0.844″", "열 하나의 폭"),
        ("0.2″", "거터"), ("6", "본문 구성 패턴")]
for i, (num, label) in enumerate(kpis):
    x = ZONE["margin"] + i * (w4 + GUTTER)
    card(s, x, ZONE["body_top"], w4, 1.6, line_color=BRAND["border"])
    textbox(s, x + 0.25, ZONE["body_top"] + 0.30, w4 - 0.5, 0.7,
            num, 32, bold=True, color=BRAND["primary"], line=1.10)
    textbox(s, x + 0.25, ZONE["body_top"] + 1.05, w4 - 0.5, 0.4,
            label, 11, color=BRAND["text_sub"])
cd = CategoryChartData()
cd.categories = ["패턴 A", "패턴 B", "패턴 D", "패턴 F"]
cd.add_series("사용 장수", (4, 3, 2, 1))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                        Inches(ZONE["margin"]), Inches(4.30),
                        Inches(CONTENT_W), Inches(2.35), cd)
ch = gf.chart
ch.has_legend = False
ch.plots[0].series[0].format.fill.solid()
ch.plots[0].series[0].format.fill.fore_color.rgb = BRAND["primary_light"]
for ax in (ch.category_axis, ch.value_axis):
    ax.tick_labels.font.size = Pt(10)
    ax.tick_labels.font.name = BRAND["font"]
    ax.tick_labels.font.color.rgb = BRAND["text_sub"]

# ── 3. 패턴 B — 2단 비교 ──
s = prs.slides.add_slide(BLANK); bg(s, BRAND["white"])
frame(s, CHAPTER, 3, TOTAL)
head(s, "패턴 B — 두 단으로 나눠 비교", "각 단 5.97″, 거터 0.4″ — 합이 정확히 12.333″ 입니다")
colw = 5.97
for i, (t1, body) in enumerate([
    ("규격이 정하는 것",
     "· 다섯 존의 Y 좌표\n· 서체 하나와 굵기 척도\n· 브랜드 색과 그라디언트 예산\n· 12열 그리드와 거터"),
    ("그때그때 정하는 것",
     "· 본문 상자 안의 구성\n· 어떤 차트를 쓸 것인가\n· 카드 몇 장으로 나될 것인가\n· 무엇을 강조할 것인가")]):
    x = ZONE["margin"] + i * (colw + 0.4)
    textbox(s, x, ZONE["body_top"], colw, 0.4, t1, 18, bold=True)
    textbox(s, x, ZONE["body_top"] + 0.55, colw, 2.0, body, 13,
            color=BRAND["text_main"], line=1.5)
card(s, ZONE["margin"], 5.55, CONTENT_W, 1.0, fill=BRAND["divider"])
textbox(s, ZONE["margin"] + 0.3, 5.80, CONTENT_W - 0.6, 0.5,
        "그래서 무엇인가 — 고정된 틀이 있어야 본문에 쓸 판단력이 남습니다.",
        14, bold=True)

# ── 4. 패턴 D — 프로세스 흐름 + 금액 표기 예시 ──
s = prs.slides.add_slide(BLANK); bg(s, BRAND["white"])
frame(s, CHAPTER, 4, TOTAL, "가상 데이터 · 서식 예시용")
head(s, "패턴 D — 만드는 순서", "금액은 줄이지 않고 그대로 씁니다")
steps = ["내용 확정", "패턴 선택", "좌표 배치", "숫자 검사", "렌더 확인"]
sw = kpi_widths(5)
for i, st in enumerate(steps):
    x = ZONE["margin"] + i * (sw + GUTTER)
    card(s, x, ZONE["body_top"], sw, 1.5, line_color=BRAND["border"])
    textbox(s, x + 0.2, ZONE["body_top"] + 0.25, sw - 0.4, 0.5,
            f"{i+1}", 22, bold=True, color=BRAND["primary"])
    textbox(s, x + 0.2, ZONE["body_top"] + 0.85, sw - 0.4, 0.4, st, 13, bold=True)
textbox(s, ZONE["margin"], 4.35, CONTENT_W, 0.4, "금액 표기 예시", 18, bold=True)
rows = [("올바른 표기", won(1234567890)), ("올바른 표기", won(358000000)),
        ("쓰지 않는 표기", "12.3억 · 123,456만 · 3억원 미만")]
for i, (k, v) in enumerate(rows):
    y = 4.90 + i * 0.5
    textbox(s, ZONE["margin"], y, 3.0, 0.4, k, 13, color=BRAND["text_sub"])
    textbox(s, ZONE["margin"] + 3.2, y, 8.0, 0.4, v, 13, bold=(i < 2),
            color=BRAND["text_main"] if i < 2 else BRAND["text_muted"])

# ── 5. 마감 — 어두운 면 ──
s = prs.slides.add_slide(BLANK); bg(s, BRAND["dark_surface"])
frame(s, CHAPTER, 5, TOTAL, dark=True)
head(s, "브랜드는 한 파일에서만 바꿉니다",
     "brand.config.md 를 고치면 이 다섯 장의 색과 서체가 함께 바뀝니다", dark=True)
textbox(s, ZONE["margin"], ZONE["body_top"] + 0.6, CONTENT_W, 2.2,
        "1.  brand.config.md 에서 색·서체·로고·통화를 바꾼다\n"
        "2.  assets/logo-placeholder.png 를 같은 이름으로 덮어쓴다\n"
        "3.  design-system.md 를 프로젝트 지침에 붙여넣는다\n"
        "4.  스킬 3개를 올린다",
        18, color=BRAND["white"], line=1.9)

out = os.path.join(os.path.dirname(__file__), "sample-deck.pptx")
prs.save(out)
print("저장:", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "장")
